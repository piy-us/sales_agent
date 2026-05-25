"""
parser.py
─────────────────────────────────────────────────────────────────────────────
STEP 1  Parse PPTX with Docling (text / tables) and python-pptx (visuals,
        metadata, speaker-notes).
STEP 2  Describe every visual with Azure OpenAI vision (GPT-4o via Foundry).
STEP 4  Assemble overlapping slide-window chunks:
          window size  = 4 slides
          overlap      = 2 slides  (i.e. stride = window - overlap = 2)

pip install docling python-pptx openai azure-storage-blob
"""

from __future__ import annotations

import base64
import hashlib
import os
from io import BytesIO
from pathlib import Path
from typing import Optional

from docling.document_converter import DocumentConverter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openai import AzureOpenAI

# ── Azure OpenAI / Foundry config ─────────────────────────────────────────
AZURE_FOUNDRY_ENDPOINT = os.environ.get(
    "AZURE_FOUNDRY_ENDPOINT", "https://services.ai.azure.com/v1"
)
AZURE_OPENAI_KEY       = os.environ["AZURE_OPENAI_KEY"]
AZURE_OPENAI_API_VER   = os.environ.get("AZURE_OPENAI_API_VERSION", "2025-01-01-preview")

# Deployment names – set these to whatever you deployed in AI Foundry
VISION_DEPLOYMENT      = os.environ.get("AZURE_VISION_DEPLOYMENT", "gpt-4o")

# Chunking strategy
CHUNK_WINDOW   = 4   # slides per chunk
CHUNK_OVERLAP  = 2   # overlapping slides between consecutive chunks
CHUNK_STRIDE   = CHUNK_WINDOW - CHUNK_OVERLAP   # = 2


# ── Vision client (Azure OpenAI via Foundry endpoint) ─────────────────────
_vision_client: Optional[AzureOpenAI] = None


def _get_vision_client() -> AzureOpenAI:
    global _vision_client
    if _vision_client is None:
        _vision_client = AzureOpenAI(
            azure_endpoint = AZURE_FOUNDRY_ENDPOINT,
            api_key        = AZURE_OPENAI_KEY,
            api_version    = AZURE_OPENAI_API_VER,
        )
    return _vision_client


# ═══════════════════════════════════════════════════════════════════
# STEP 1 — Parse
# ═══════════════════════════════════════════════════════════════════

def parse_pptx(path: str) -> dict:
    """
    Returns:
        {
          "file_metadata": {...},
          "slides": {
              1: {"title": ..., "notes": ..., "text_md": ..., "visuals": [...]},
              ...
          }
        }
    """
    print(f"  Parsing {path} …")
    file = Path(path)

    # ── Docling for text + tables ──────────────────────────────────
    converter = DocumentConverter()
    result    = converter.convert(path)
    doc       = result.document
    full_md   = doc.export_to_markdown()

    # Split markdown by slide (Docling uses "## Slide N" / "# Slide N")
    slide_md: dict[int, str] = {}
    current_slide = 0
    current_lines: list[str] = []

    for line in full_md.splitlines():
        if line.startswith("## Slide ") or line.startswith("# Slide "):
            if current_lines:
                slide_md[current_slide] = "\n".join(current_lines).strip()
            try:
                current_slide = int(line.split()[-1])
            except ValueError:
                current_slide += 1
            current_lines = [line]
        else:
            current_lines.append(line)
    if current_lines:
        slide_md[current_slide] = "\n".join(current_lines).strip()

    # ── python-pptx for images, charts, metadata ──────────────────
    prs = Presentation(path)
    cp  = prs.core_properties

    file_metadata = {
        "source_file":   file.name,
        "author":        cp.author or "",
        "modified_date": str(cp.modified or ""),
        "slide_count":   len(prs.slides),
    }

    slides_data: dict[int, dict] = {}
    for slide_num, slide in enumerate(prs.slides, start=1):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.lower().startswith("title"):
                title = shape.text_frame.text.strip()
                break

        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        visuals: list[dict] = []
        for shape in slide.shapes:
            if shape.has_chart:
                blob = _get_chart_blob(shape)
                if blob:
                    visuals.append({
                        "shape_name":  shape.name,
                        "blob":        blob,
                        "is_chart":    True,
                        "slide_title": title,
                    })
                else:
                    visuals.append({
                        "shape_name":   shape.name,
                        "chart_series": _chart_series_text(shape),
                        "is_chart":     True,
                        "is_data_only": True,
                        "slide_title":  title,
                    })
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                visuals.append({
                    "shape_name":  shape.name,
                    "blob":        shape.image.blob,
                    "is_chart":    False,
                    "slide_title": title,
                })

        slides_data[slide_num] = {
            "title":   title,
            "notes":   notes,
            "text_md": slide_md.get(slide_num, ""),
            "visuals": visuals,
        }

    return {"file_metadata": file_metadata, "slides": slides_data}


def _get_chart_blob(shape) -> bytes | None:
    """Try to get the cached PNG preview that PowerPoint embeds for charts."""
    try:
        chart_part = shape.chart.part
        for rel in chart_part.rels.values():
            if "image" in rel.reltype:
                return rel.target_part.blob
    except Exception:
        pass
    return None


def _chart_series_text(shape) -> str:
    parts = []
    try:
        for s in shape.chart.series:
            vals = [v for v in s.values if v is not None]
            parts.append(f"{s.name}: {vals}")
    except Exception:
        pass
    return "; ".join(parts)


# ═══════════════════════════════════════════════════════════════════
# STEP 2 — Visual description via Azure OpenAI GPT-4o vision
# ═══════════════════════════════════════════════════════════════════

CHART_PROMPT = """You are describing a chart from a business presentation for a RAG system.
Slide title: {title}

Provide a structured description covering:
1. Chart type (bar, line, pie, scatter, etc.)
2. What is being measured (metric, KPI, units)
3. All axis labels and legend entries
4. Every visible data point or value (be precise)
5. The main trend or insight in one sentence
6. Any annotations, callouts, or highlighted values

Be exhaustive — someone will use this to answer data questions."""

IMAGE_PROMPT = """You are describing an image from a business presentation for a RAG system.
Slide title: {title}

Provide a structured description covering:
1. What type of visual this is (diagram, photo, screenshot, architecture, etc.)
2. All visible text, labels, and annotations
3. All components and how they connect or relate
4. The key information or message conveyed
5. Any data, metrics, or specific values visible

Be exhaustive — someone will use this to answer questions about this content."""


def describe_visual(blob: bytes, slide_title: str, is_chart: bool) -> str:
    """Call GPT-4o (vision) via Azure AI Foundry endpoint to describe a visual."""
    prompt_template = CHART_PROMPT if is_chart else IMAGE_PROMPT
    prompt = prompt_template.format(title=slide_title or "Unknown")

    # Detect media type
    if blob[:3] == b"\xff\xd8\xff":
        media_type = "image/jpeg"
    elif blob[:4] == b"\x89PNG":
        media_type = "image/png"
    else:
        media_type = "image/png"

    b64_image = base64.standard_b64encode(blob).decode()
    data_url  = f"data:{media_type};base64,{b64_image}"

    try:
        client   = _get_vision_client()
        response = client.chat.completions.create(
            model    = VISION_DEPLOYMENT,
            max_tokens = 600,
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type":      "image_url",
                            "image_url": {"url": data_url, "detail": "high"},
                        },
                        {"type": "text", "text": prompt},
                    ],
                }
            ],
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"[Description unavailable: {e}]"


# ═══════════════════════════════════════════════════════════════════
# STEP 4 — Assemble overlapping slide-window chunks
#
#   Strategy:
#     • Each TEXT chunk spans CHUNK_WINDOW (=4) consecutive slides.
#     • Stride = CHUNK_WINDOW - CHUNK_OVERLAP = 2, so adjacent chunks
#       share CHUNK_OVERLAP (=2) slides — giving rich contextual overlap
#       for slides that are too small on their own.
#     • Visual chunks (charts / images) remain per-slide so the image_uri
#       mapping stays unambiguous.
# ═══════════════════════════════════════════════════════════════════

def assemble_chunks(parsed: dict, upload_fn=None) -> list[dict]:
    """
    Parameters
    ----------
    parsed      : output of parse_pptx()
    upload_fn   : callable(blob_bytes, filename) -> str  — optional; if
                  supplied, visuals are uploaded and the returned URI is
                  stored on the chunk.  Pass uploader.upload_image_to_blob.

    Returns
    -------
    List of chunk dicts ready for embedding + indexing.
    """
    file_meta   = parsed["file_metadata"]
    slides_data = parsed["slides"]
    slide_nums  = sorted(slides_data.keys())   # [1, 2, 3, …, N]
    chunks: list[dict] = []

    # ── sliding-window TEXT chunks ────────────────────────────────
    start = 0
    while start < len(slide_nums):
        window = slide_nums[start : start + CHUNK_WINDOW]   # up to 4 slides
        window_texts: list[str] = []

        for sn in window:
            slide = slides_data[sn]
            parts: list[str] = []
            if slide["title"]:
                parts.append(f"### Slide {sn}: {slide['title']}")
            if slide["text_md"]:
                parts.append(slide["text_md"])
            if slide["notes"]:
                parts.append(f"[Speaker notes]: {slide['notes']}")
            if parts:
                window_texts.append("\n".join(parts))

        text_body = "\n\n---\n\n".join(window_texts).strip()
        if text_body:
            first_slide = window[0]
            last_slide  = window[-1]
            chunk_id    = _chunk_id(
                file_meta["source_file"],
                first_slide,
                f"text_w{first_slide}_{last_slide}",
            )
            chunks.append({
                "id":           chunk_id,
                "text":         text_body,
                "content_type": "text",
                "image_uri":    None,
                "metadata": {
                    **file_meta,
                    "slide_number":      first_slide,
                    "slide_number_end":  last_slide,
                    "slide_title":       slides_data[first_slide]["title"]
                                         or f"Slide {first_slide}",
                    "content_type":      "text",
                    "chunk_window":      list(window),
                },
            })

        start += CHUNK_STRIDE
        if start < len(slide_nums) and start + CHUNK_WINDOW > len(slide_nums):
            # ensure last partial window is always included
            start = max(len(slide_nums) - CHUNK_WINDOW, start)

    # ── per-slide VISUAL chunks ───────────────────────────────────
    for slide_num in slide_nums:
        slide     = slides_data[slide_num]
        base_meta = {
            **file_meta,
            "slide_number": slide_num,
            "slide_title":  slide["title"] or f"Slide {slide_num}",
        }

        for i, visual in enumerate(slide["visuals"]):
            # Chart with data only (no PNG preview available)
            if visual.get("is_data_only"):
                desc = (
                    f"[Chart data — {visual['shape_name']}]\n"
                    f"{visual['chart_series']}"
                )
                chunks.append({
                    "id":           _chunk_id(file_meta["source_file"], slide_num, f"chart_{i}"),
                    "text":         f"{slide['title']}\n\n{desc}",
                    "content_type": "chart",
                    "image_uri":    None,
                    "metadata":     {**base_meta, "content_type": "chart"},
                })
                continue

            blob = visual["blob"]

            # Upload image if an uploader was provided
            image_uri: str | None = None
            if upload_fn is not None:
                img_filename = _chunk_id(
                    file_meta["source_file"], slide_num, f"img_{i}"
                ) + ".png"
                image_uri = upload_fn(blob, img_filename)

            # Describe with vision LLM
            visual_type = "chart" if visual["is_chart"] else "image"
            print(
                f"    🔍 Describing {visual_type}: "
                f"{visual['shape_name']} (slide {slide_num})"
            )
            description = describe_visual(blob, slide["title"], visual["is_chart"])

            chunk_text = (
                f"{slide['title']}\n\n"
                f"[{visual_type.upper()} — {visual['shape_name']}]\n"
                f"{description}"
            )
            chunks.append({
                "id":           _chunk_id(
                    file_meta["source_file"], slide_num, f"{visual_type}_{i}"
                ),
                "text":         chunk_text,
                "content_type": visual_type,
                "image_uri":    image_uri,
                "metadata": {
                    **base_meta,
                    "content_type": visual_type,
                    "image_uri":    image_uri,
                    "shape_name":   visual["shape_name"],
                },
            })

    return chunks


def _chunk_id(filename: str, slide: int, suffix: str) -> str:
    raw = f"{filename}_slide{slide}_{suffix}"
    return hashlib.md5(raw.encode()).hexdigest()
