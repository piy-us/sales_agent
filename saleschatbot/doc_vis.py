"""
extract_images.py
Extracts all images and charts from a PPTX file and saves them to a folder.

pip install python-pptx Pillow
python extract_images.py your_file.pptx
"""

import sys, os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images(pptx_path: str):
    prs      = Presentation(pptx_path)
    out_dir  = Path(pptx_path).stem + "_images"
    os.makedirs(out_dir, exist_ok=True)

    saved = []

    for slide_num, slide in enumerate(prs.slides, start=1):

        # get slide title for the filename
        title = f"slide{slide_num}"
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.lower().startswith("title"):
                t = shape.text_frame.text.strip()
                if t:
                    # sanitize for filename
                    safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in t)
                    title = f"slide{slide_num}_{safe[:40]}"
                break

        for shape_idx, shape in enumerate(slide.shapes):

            # ── embedded picture ──────────────────────────────────
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob        = shape.image.blob
                ext         = shape.image.content_type.split("/")[-1]  # png, jpeg, etc.
                ext         = ext.replace("jpeg", "jpg")
                filename    = f"{title}_img{shape_idx}.{ext}"
                out_path    = Path(out_dir) / filename

                with open(out_path, "wb") as f:
                    f.write(blob)

                saved.append((slide_num, "image", shape.name, str(out_path)))
                print(f"  ✅ Slide {slide_num} | image  | {shape.name:<30} → {filename}")

            # ── native chart object ───────────────────────────────
            elif shape.has_chart:
                blob = None

                # PowerPoint caches a PNG preview inside every chart — grab it
                try:
                    chart_part = shape.chart.part
                    for rel in chart_part.rels.values():
                        if "image" in rel.reltype:
                            blob = rel.target_part.blob
                            break
                except Exception:
                    pass

                if blob:
                    filename = f"{title}_chart{shape_idx}.png"
                    out_path = Path(out_dir) / filename
                    with open(out_path, "wb") as f:
                        f.write(blob)
                    saved.append((slide_num, "chart", shape.name, str(out_path)))
                    print(f"  ✅ Slide {slide_num} | chart  | {shape.name:<30} → {filename}")
                else:
                    print(f"  ⚠️  Slide {slide_num} | chart  | {shape.name:<30} → no cached PNG found")

    print(f"\n  Saved {len(saved)} files to ./{out_dir}/")
    return saved


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else "sample_presentation.pptx"
    if not Path(path).exists():
        print(f"File not found: {path}")
        sys.exit(1)

    print(f"\nExtracting images from: {path}\n")
    extract_images(path)