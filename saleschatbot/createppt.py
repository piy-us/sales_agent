"""
Creates a rich sample PPTX file to stress-test both parsers.
Includes: text, tables, images, charts, speaker notes, multi-level bullets, shapes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import io, math

# ── tiny generated image (red gradient PNG, no external deps) ────────────────
def _make_png_bytes():
    import zlib, struct
    w, h = 60, 40
    raw = b""
    for y in range(h):
        raw += b"\x00"
        for x in range(w):
            r = int(255 * x / w)
            g = int(180 * y / h)
            b = 80
            raw += bytes([r, g, b])
    def chunk(tag, data):
        c = zlib.crc32(tag + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", c)
    png = (b"\x89PNG\r\n\x1a\n"
           + chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
           + chunk(b"IDAT", zlib.compress(raw))
           + chunk(b"IEND", b""))
    return png


def create_sample_pptx(path: str = "sample_presentation.pptx"):
    prs = Presentation()
    W, H = prs.slide_width, prs.slide_height

    # ── helper ───────────────────────────────────────────────────────────────
    def add_textbox(slide, text, left, top, width, height,
                    bold=False, size=18, color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = RGBColor(*color)
        return tb

    # ════════════════════════════════════════════════════════════════════════
    # Slide 1 – Title slide
    # ════════════════════════════════════════════════════════════════════════
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Enterprise RAG on PowerPoint Data"
    slide.placeholders[1].text = (
        "A Deep Dive into PPT Parsing Capabilities\n"
        "Docling vs Unstructured vs python-pptx"
    )
    slide.notes_slide.notes_text_frame.text = (
        "Welcome slide. Key message: choosing the right parser is critical "
        "for RAG quality. We will benchmark three approaches today."
    )

    # ════════════════════════════════════════════════════════════════════════
    # Slide 2 – Bullet hierarchy
    # ════════════════════════════════════════════════════════════════════════
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Why PPT Parsing Is Hard"
    tf = slide.placeholders[1].text_frame
    tf.text = "Structural complexity"
    tf.paragraphs[0].level = 0

    items = [
        (1, "Slides carry semantic hierarchy (title → body → sub-bullet)"),
        (2, "Formatting encodes meaning (bold = key term, red = warning)"),
        (1, "Mixed content types in one file"),
        (2, "Text, tables, images, SmartArt, charts, embedded objects"),
        (2, "Speaker notes as supplementary context"),
        (0, "Metadata matters for enterprise RAG"),
        (1, "Slide number, author, last-modified date"),
        (1, "Section names and slide titles as chunk boundaries"),
    ]
    for level, text in items:
        p = tf.add_paragraph()
        p.text = text
        p.level = level

    slide.notes_slide.notes_text_frame.text = (
        "Emphasize that naive text extraction loses the hierarchy. "
        "A level-2 bullet under 'Tables' means something different from "
        "the same text under 'Images'."
    )

    # ════════════════════════════════════════════════════════════════════════
    # Slide 3 – Table
    # ════════════════════════════════════════════════════════════════════════
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Parser Feature Comparison"

    rows, cols = 5, 4
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.5), Inches(9), Inches(3.5)
    ).table

    headers = ["Feature", "python-pptx", "Unstructured", "Docling"]
    widths  = [Inches(3), Inches(2), Inches(2), Inches(2)]
    for i, (h, w) in enumerate(zip(headers, widths)):
        table.columns[i].width = w
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    data = [
        ["Text extraction",    "✅ Native",   "✅ Yes",    "✅ Yes"],
        ["Table extraction",   "✅ Native",   "✅ Yes",    "✅ Yes"],
        ["Image extraction",   "⚠️  Manual",  "✅ Yes",    "✅ Yes"],
        ["Markdown export",    "❌ No",       "✅ Yes",    "✅ Yes"],
    ]
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    slide.notes_slide.notes_text_frame.text = (
        "This table is the core comparison slide. "
        "Docling and Unstructured both wrap python-pptx internally "
        "but add OCR and layout intelligence on top."
    )

    # ════════════════════════════════════════════════════════════════════════
    # Slide 4 – Chart
    # ════════════════════════════════════════════════════════════════════════
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Benchmark: Parsing Speed (seconds)"

    cd = ChartData()
    cd.categories = ["10 slides", "50 slides", "200 slides"]
    cd.add_series("python-pptx",  (0.3,  1.2,   5.1))
    cd.add_series("Unstructured", (1.8,  8.4,  33.0))
    cd.add_series("Docling",      (4.2, 19.7,  78.5))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(8), Inches(4.5),
        cd
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Lower is faster"
    chart.has_legend = True

    slide.notes_slide.notes_text_frame.text = (
        "Synthetic benchmark — actual numbers depend on hardware and "
        "whether GPU-accelerated OCR is enabled. "
        "Docling's slowness is offset by much richer structural output."
    )

    # ════════════════════════════════════════════════════════════════════════
    # Slide 5 – Embedded image + shapes
    # ════════════════════════════════════════════════════════════════════════
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Architecture: RAG Pipeline"

    png_bytes = _make_png_bytes()
    img_stream = io.BytesIO(png_bytes)
    slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.4), Inches(2.5), Inches(1.8))

    boxes = [
        (Inches(3.0), "1. Ingest\nPPTX files", RGBColor(0x1F, 0x77, 0xB4)),
        (Inches(5.0), "2. Parse\n& Chunk",     RGBColor(0xFF, 0x7F, 0x0E)),
        (Inches(7.0), "3. Embed\n& Index",     RGBColor(0x2C, 0xA0, 0x2C)),
    ]
    for left, label, color in boxes:
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, Inches(2.0), Inches(1.8), Inches(1.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf = shape.text_frame
        tf.text = label
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(12)

    add_textbox(slide,
        "Each box = one stage. The parser determines chunk quality.",
        Inches(0.3), Inches(4.0), Inches(9.3), Inches(0.8),
        size=13, color=(80, 80, 80))

    slide.notes_slide.notes_text_frame.text = (
        "The image on the left is a programmatically generated gradient — "
        "parsers should detect it as an embedded image element. "
        "The coloured boxes are native PPTX shapes."
    )

    # ════════════════════════════════════════════════════════════════════════
    # Slide 6 – Rich text formatting
    # ════════════════════════════════════════════════════════════════════════
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Formatting as Signal"

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    styles = [
        ("Normal paragraph — baseline text parsers must capture.\n", False, False, None, 14),
        ("Bold signals key terms or headings within body text.\n",    True,  False, None, 14),
        ("Italic often marks technical terms or citations.\n",        False, True,  None, 14),
        ("Red text typically marks warnings or critical notes.\n",    False, False, RGBColor(0xC0, 0x00, 0x00), 14),
        ("Large text (20 pt) signals section headers inside slides.\n", False, True, RGBColor(0x1F, 0x49, 0x7D), 20),
        ("Hyperlink text: https://docling.github.io",                False, False, RGBColor(0x00, 0x56, 0xB3), 14),
    ]

    first = True
    for text, bold, italic, color, size in styles:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.italic = italic
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = color

    slide.notes_slide.notes_text_frame.text = (
        "Rich formatting is crucial for RAG chunking heuristics. "
        "A good parser returns font size, bold, italic, and color — "
        "not just raw string content."
    )

    # ════════════════════════════════════════════════════════════════════════
    # Slide 7 – Conclusion
    # ════════════════════════════════════════════════════════════════════════
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Recommendation"
    slide.placeholders[1].text = (
        "Use Docling for production RAG (best structure + Markdown output)\n"
        "Use Unstructured for speed-sensitive pipelines\n"
        "Use python-pptx directly for metadata + notes extraction"
    )
    slide.notes_slide.notes_text_frame.text = (
        "Final recommendation: hybrid approach works best. "
        "python-pptx for metadata pass, Docling for content pass."
    )

    prs.save(path)
    print(f"✅  Sample PPTX created → {path}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    create_sample_pptx()