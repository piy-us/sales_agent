"""
enterprise_pptx_ingestion.py
Full ingestion pipeline for Azure enterprise RAG.
Runs completely offline after first model download.

pip install docling azure-search-documents azure-storage-blob anthropic python-pptx
"""

import os, json, base64, hashlib
from pathlib import Path
from io import BytesIO

from docling.document_converter import DocumentConverter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import anthropic
from azure.storage.blob import BlobServiceClient
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from azure.search.documents.indexes.models import (
    SearchIndex, SimpleField, SearchableField,
    SearchFieldDataType, VectorSearch,
    HnswAlgorithmConfiguration, VectorSearchProfile,
    SearchField
)
from azure.core.credentials import AzureKeyCredential

# ── config (use environment variables or Azure Key Vault in production) ──────
AZURE_SEARCH_ENDPOINT  = os.environ["AZURE_SEARCH_ENDPOINT"]
AZURE_SEARCH_KEY       = os.environ["AZURE_SEARCH_KEY"]
AZURE_STORAGE_CONN     = os.environ["AZURE_STORAGE_CONN"]
ANTHROPIC_KEY          = os.environ["ANTHROPIC_KEY"]
INDEX_NAME             = "pptx-rag-index"
BLOB_CONTAINER         = "slide-images"


# ═══════════════════════════════════════════════════════════════════
# STEP 1 — Parse PPTX with Docling (text, tables, structure)
#          AND python-pptx (images, charts, metadata)
# ═══════════════════════════════════════════════════════════════════

def parse_pptx(path: str) -> dict:
    """
    Returns a dict keyed by slide_number containing:
      - text_md:    Docling markdown for that slide's text/tables
      - images:     list of {shape_name, blob, is_chart, slide_title}
      - metadata:   file-level properties
    """
    print(f"  Parsing {path} …")
    file = Path(path)

    # ── Docling for text + tables ────────────────────────────────
    converter = DocumentConverter()
    result    = converter.convert(path)
    doc       = result.document
    full_md   = doc.export_to_markdown()

    # Split markdown by slide (Docling uses "## Slide N" headings for PPTX)
    slide_md = {}
    current_slide = 0
    current_lines = []
    for line in full_md.splitlines():
        if line.startswith("## Slide ") or line.startswith("# Slide "):
            if current_lines:
                slide_md[current_slide] = "\n".join(current_lines).strip()
            try:
                current_slide = int(line.split()[-1])
            except ValueError:
                current_slide += 1
            current_lines = [line]
        else:
            current_lines.append(line)
    if current_lines:
        slide_md[current_slide] = "\n".join(current_lines).strip()

    # ── python-pptx for images, charts, metadata ─────────────────
    prs = Presentation(path)
    cp  = prs.core_properties

    file_metadata = {
        "source_file":   file.name,
        "author":        cp.author or "",
        "modified_date": str(cp.modified or ""),
        "slide_count":   len(prs.slides),
    }

    slides_data = {}
    for slide_num, slide in enumerate(prs.slides, start=1):
        # get slide title
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.lower().startswith("title"):
                title = shape.text_frame.text.strip()
                break

        # get speaker notes
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        # extract images and charts
        visuals = []
        for shape in slide.shapes:
            # native chart object
            if shape.has_chart:
                blob = _get_chart_blob(shape)
                if blob:
                    visuals.append({
                        "shape_name":  shape.name,
                        "blob":        blob,
                        "is_chart":    True,
                        "slide_title": title,
                    })
                else:
                    # fallback: serialize chart data as text
                    visuals.append({
                        "shape_name":   shape.name,
                        "chart_series": _chart_series_text(shape),
                        "is_chart":     True,
                        "is_data_only": True,
                        "slide_title":  title,
                    })

            # embedded picture
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                visuals.append({
                    "shape_name":  shape.name,
                    "blob":        shape.image.blob,
                    "is_chart":    False,
                    "slide_title": title,
                })

        slides_data[slide_num] = {
            "title":    title,
            "notes":    notes,
            "text_md":  slide_md.get(slide_num, ""),
            "visuals":  visuals,
        }

    return {"file_metadata": file_metadata, "slides": slides_data}


def _get_chart_blob(shape) -> bytes | None:
    """Try to get the cached PNG preview that PowerPoint embeds for charts."""
    try:
        chart_part = shape.chart.part
        for rel in chart_part.rels.values():
            if "image" in rel.reltype:
                return rel.target_part.blob
    except Exception:
        pass
    return None


def _chart_series_text(shape) -> str:
    parts = []
    try:
        for s in shape.chart.series:
            vals = [v for v in s.values if v is not None]
            parts.append(f"{s.name}: {vals}")
    except Exception:
        pass
    return "; ".join(parts)


# ═══════════════════════════════════════════════════════════════════
# STEP 2 — Describe every visual with Vision LLM (once at ingest)
# ═══════════════════════════════════════════════════════════════════

vision_client = anthropic.Anthropic(api_key=ANTHROPIC_KEY)

CHART_PROMPT = """You are describing a chart from a business presentation for a RAG system.
Slide title: {title}

Provide a structured description covering:
1. Chart type (bar, line, pie, scatter, etc.)
2. What is being measured (metric, KPI, units)
3. All axis labels and legend entries
4. Every visible data point or value (be precise)
5. The main trend or insight in one sentence
6. Any annotations, callouts, or highlighted values

Be exhaustive — someone will use this to answer data questions."""

IMAGE_PROMPT = """You are describing an image from a business presentation for a RAG system.
Slide title: {title}

Provide a structured description covering:
1. What type of visual this is (diagram, photo, screenshot, architecture, etc.)
2. All visible text, labels, and annotations
3. All components and how they connect or relate
4. The key information or message conveyed
5. Any data, metrics, or specific values visible

Be exhaustive — someone will use this to answer questions about this content."""


def describe_visual(blob: bytes, slide_title: str, is_chart: bool) -> str:
    prompt_template = CHART_PROMPT if is_chart else IMAGE_PROMPT
    prompt = prompt_template.format(title=slide_title or "Unknown")

    media_type = "image/png"
    if blob[:3] == b'\xff\xd8\xff':
        media_type = "image/jpeg"

    try:
        response = vision_client.messages.create(
            model      = "claude-sonnet-4-20250514",
            max_tokens = 600,
            messages   = [{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type":       "base64",
                            "media_type": media_type,
                            "data":       base64.standard_b64encode(blob).decode()
                        }
                    },
                    {"type": "text", "text": prompt}
                ]
            }]
        )
        return response.content[0].text.strip()
    except Exception as e:
        return f"[Description unavailable: {e}]"


# ═══════════════════════════════════════════════════════════════════
# STEP 3 — Upload images to Azure Blob, get back SAS URLs
# ═══════════════════════════════════════════════════════════════════

blob_service = BlobServiceClient.from_connection_string(AZURE_STORAGE_CONN)

def upload_image_to_blob(blob_bytes: bytes, filename: str) -> str:
    """Upload PNG to Azure Blob and return the URL."""
    container_client = blob_service.get_container_client(BLOB_CONTAINER)
    try:
        container_client.create_container()
    except Exception:
        pass  # already exists

    blob_client = container_client.get_blob_client(filename)
    blob_client.upload_blob(blob_bytes, overwrite=True,
                            content_settings={"content_type": "image/png"})
    return blob_client.url   # use SAS URL in production


# ═══════════════════════════════════════════════════════════════════
# STEP 4 — Assemble chunks
# ═══════════════════════════════════════════════════════════════════

def assemble_chunks(parsed: dict) -> list[dict]:
    """
    One text chunk per slide.
    One additional chunk per visual (description embedded in text).
    """
    chunks   = []
    file_meta = parsed["file_metadata"]

    for slide_num, slide in parsed["slides"].items():
        base_meta = {
            **file_meta,
            "slide_number": slide_num,
            "slide_title":  slide["title"] or f"Slide {slide_num}",
        }

        # ── text chunk ───────────────────────────────────────────
        text_parts = []
        if slide["title"]:
            text_parts.append(slide["title"])
        if slide["text_md"]:
            text_parts.append(slide["text_md"])
        if slide["notes"]:
            text_parts.append(f"[Speaker notes]: {slide['notes']}")

        text_body = "\n\n".join(text_parts).strip()
        if text_body:
            chunks.append({
                "id":           _chunk_id(file_meta["source_file"], slide_num, "text"),
                "text":         text_body,
                "content_type": "text",
                "image_uri":    None,
                "metadata":     {**base_meta, "content_type": "text"},
            })

        # ── visual chunks ─────────────────────────────────────────
        for i, visual in enumerate(slide["visuals"]):

            # chart with only data (no PNG)
            if visual.get("is_data_only"):
                desc = f"[Chart data — {visual['shape_name']}]\n{visual['chart_series']}"
                chunks.append({
                    "id":           _chunk_id(file_meta["source_file"], slide_num, f"chart_{i}"),
                    "text":         f"{slide['title']}\n\n{desc}",
                    "content_type": "chart",
                    "image_uri":    None,
                    "metadata":     {**base_meta, "content_type": "chart"},
                })
                continue

            blob = visual["blob"]

            # upload to blob storage
            img_filename = _chunk_id(file_meta["source_file"], slide_num, f"img_{i}") + ".png"
            image_uri    = upload_image_to_blob(blob, img_filename)

            # describe with vision LLM
            print(f"    🔍 Describing {'chart' if visual['is_chart'] else 'image'}: {visual['shape_name']}")
            description = describe_visual(blob, slide["title"], visual["is_chart"])

            visual_type = "chart" if visual["is_chart"] else "image"

            # description goes IN the text so it gets embedded
            chunk_text = (
                f"{slide['title']}\n\n"
                f"[{visual_type.upper()} — {visual['shape_name']}]\n"
                f"{description}"
            )

            chunks.append({
                "id":           _chunk_id(file_meta["source_file"], slide_num, f"{visual_type}_{i}"),
                "text":         chunk_text,
                "content_type": visual_type,
                "image_uri":    image_uri,   # stored in metadata for tool fetch
                "metadata":     {
                    **base_meta,
                    "content_type": visual_type,
                    "image_uri":    image_uri,
                    "shape_name":   visual["shape_name"],
                },
            })

    return chunks


def _chunk_id(filename: str, slide: int, suffix: str) -> str:
    raw = f"{filename}_slide{slide}_{suffix}"
    return hashlib.md5(raw.encode()).hexdigest()


# ═══════════════════════════════════════════════════════════════════
# STEP 5 — Embed + index into Azure AI Search
# ═══════════════════════════════════════════════════════════════════

def get_embedding(text: str) -> list[float]:
    """Use any embedding model — text-embedding-3-small shown here."""
    import openai
    client = openai.AzureOpenAI(
        azure_endpoint = os.environ["AZURE_OPENAI_ENDPOINT"],
        api_key        = os.environ["AZURE_OPENAI_KEY"],
        api_version    = "2024-02-01"
    )
    response = client.embeddings.create(
        input = text,
        model = "text-embedding-3-small"
    )
    return response.data[0].embedding


def create_index_if_missing():
    index_client = SearchIndexClient(
        endpoint   = AZURE_SEARCH_ENDPOINT,
        credential = AzureKeyCredential(AZURE_SEARCH_KEY)
    )
    fields = [
        SimpleField(name="id",           type=SearchFieldDataType.String, key=True),
        SearchableField(name="text",     type=SearchFieldDataType.String),
        SimpleField(name="content_type", type=SearchFieldDataType.String, filterable=True),
        SimpleField(name="source_file",  type=SearchFieldDataType.String, filterable=True),
        SimpleField(name="slide_number", type=SearchFieldDataType.Int32,  filterable=True, sortable=True),
        SimpleField(name="slide_title",  type=SearchFieldDataType.String),
        SimpleField(name="image_uri",    type=SearchFieldDataType.String),
        SearchField(
            name               = "embedding",
            type               = SearchFieldDataType.Collection(SearchFieldDataType.Single),
            searchable         = True,
            vector_search_dimensions = 1536,
            vector_search_profile_name = "hnsw-profile"
        ),
    ]
    vector_search = VectorSearch(
        algorithms = [HnswAlgorithmConfiguration(name="hnsw")],
        profiles   = [VectorSearchProfile(name="hnsw-profile", algorithm_configuration_name="hnsw")]
    )
    index = SearchIndex(name=INDEX_NAME, fields=fields, vector_search=vector_search)
    index_client.create_or_update_index(index)
    print(f"  ✅ Index '{INDEX_NAME}' ready")


def index_chunks(chunks: list[dict]):
    create_index_if_missing()
    search_client = SearchClient(
        endpoint   = AZURE_SEARCH_ENDPOINT,
        index_name = INDEX_NAME,
        credential = AzureKeyCredential(AZURE_SEARCH_KEY)
    )
    documents = []
    for chunk in chunks:
        print(f"  📐 Embedding chunk: {chunk['id'][:20]}…")
        doc = {
            "id":           chunk["id"],
            "text":         chunk["text"],
            "content_type": chunk["content_type"],
            "source_file":  chunk["metadata"]["source_file"],
            "slide_number": chunk["metadata"]["slide_number"],
            "slide_title":  chunk["metadata"]["slide_title"],
            "image_uri":    chunk["image_uri"] or "",
            "embedding":    get_embedding(chunk["text"]),
        }
        documents.append(doc)

    # upload in batches of 100
    batch_size = 100
    for i in range(0, len(documents), batch_size):
        search_client.upload_documents(documents[i:i+batch_size])

    print(f"  ✅ Indexed {len(documents)} chunks")


# ═══════════════════════════════════════════════════════════════════
# STEP 6 — Query-time: retrieval + vision tool
# ═══════════════════════════════════════════════════════════════════

def retrieve(query: str, top_k: int = 5, filter_expr: str = None) -> list[dict]:
    """Hybrid search: vector + BM25 keyword."""
    from azure.search.documents.models import VectorizedQuery

    search_client = SearchClient(
        endpoint   = AZURE_SEARCH_ENDPOINT,
        index_name = INDEX_NAME,
        credential = AzureKeyCredential(AZURE_SEARCH_KEY)
    )
    vector_query = VectorizedQuery(
        vector         = get_embedding(query),
        k_nearest_neighbors = top_k,
        fields         = "embedding"
    )
    results = search_client.search(
        search_text   = query,       # BM25 keyword search
        vector_queries = [vector_query],  # vector search
        filter         = filter_expr,
        top            = top_k,
        select         = ["id", "text", "content_type", "slide_title",
                          "slide_number", "source_file", "image_uri"]
    )
    return [dict(r) for r in results]


# The vision tool definition — passed to the LLM
VISION_TOOL = {
    "name": "fetch_slide_image",
    "description": (
        "Fetch the actual image from a slide when the text description is not "
        "sufficient to answer the user's question precisely. Use this when the "
        "user asks about specific values, colors, or details in a chart or diagram."
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "image_uri": {
                "type": "string",
                "description": "The image URI from the retrieved chunk metadata"
            },
            "question": {
                "type": "string",
                "description": "The specific question to answer by looking at this image"
            }
        },
        "required": ["image_uri", "question"]
    }
}


def fetch_and_analyze_image(image_uri: str, question: str) -> str:
    """Called when the LLM invokes the vision tool."""
    import urllib.request
    with urllib.request.urlopen(image_uri) as resp:
        blob = resp.read()
    return describe_visual(blob, question, is_chart=True)


def answer_query(user_question: str) -> str:
    """Full RAG query pipeline with optional vision tool call."""
    chunks = retrieve(user_question, top_k=5)

    # build context — include image_uri in the text so LLM knows it exists
    context_parts = []
    for c in chunks:
        part = f"[Slide {c['slide_number']} — {c['slide_title']}]\n{c['text']}"
        if c.get("image_uri"):
            part += f"\n[image_uri: {c['image_uri']}]"
        context_parts.append(part)

    context = "\n\n---\n\n".join(context_parts)

    system = (
        "You are an enterprise assistant. Answer questions using the provided "
        "slide content. If a retrieved chunk mentions an image_uri and you need "
        "more detail from that image, call the fetch_slide_image tool."
    )

    messages = [{"role": "user", "content": f"Context:\n{context}\n\nQuestion: {user_question}"}]

    response = vision_client.messages.create(
        model    = "claude-sonnet-4-20250514",
        max_tokens = 1000,
        system   = system,
        tools    = [VISION_TOOL],
        messages = messages
    )

    # handle tool use
    while response.stop_reason == "tool_use":
        tool_use = next(b for b in response.content if b.type == "tool_use")
        tool_result = fetch_and_analyze_image(
            tool_use.input["image_uri"],
            tool_use.input["question"]
        )
        messages += [
            {"role": "assistant", "content": response.content},
            {"role": "user",      "content": [{
                "type":        "tool_result",
                "tool_use_id": tool_use.id,
                "content":     tool_result
            }]}
        ]
        response = vision_client.messages.create(
            model    = "claude-sonnet-4-20250514",
            max_tokens = 1000,
            system   = system,
            tools    = [VISION_TOOL],
            messages = messages
        )

    return next(b.text for b in response.content if hasattr(b, "text"))


# ═══════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import sys
    pptx_path = sys.argv[1]

    print("\n=== INGESTION ===")
    parsed = parse_pptx(pptx_path)
    chunks = assemble_chunks(parsed)
    print(f"  Assembled {len(chunks)} chunks")
    index_chunks(chunks)

    print("\n=== QUERY TEST ===")
    answer = answer_query("What does the revenue chart show for Q3?")
    print(answer)